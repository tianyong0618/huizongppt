from pptx import Presentation
from pptx.util import Inches
import os

# 创建PPT对象
prs = Presentation()

# 定义页面大小
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 截图目录
img_dir = '/Users/tianyong/skills/huizongppt/screenshots'

# 按顺序获取截图文件
img_files = [
    '1-导航页.png',
    '2-赋能中心运营概览-整体.png',
    '3-产品研发进展.png',
    '4-培训业务情况.png',
    '5-各地业务部门概览.png',
    '6-人力资源状况.png',
    '7-各地业务部门明细.png'
]

# 页面标题
slide_titles = [
    '2026一季度公司整体运营总览',
    '赋能中心运营概览',
    '产品研发进展',
    '培训业务情况',
    '各地业务部门概览',
    '人力资源状况',
    '各地业务部门明细'
]

# 为每个截图创建幻灯片
for i, (img_file, title) in enumerate(zip(img_files, slide_titles)):
    img_path = os.path.join(img_dir, img_file)
    
    if i == 0:
        # 第一页使用标题幻灯片
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
    else:
        # 其他页使用标题和内容幻灯片
        slide_layout = prs.slide_layouts[5]  # 空白幻灯片
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Inches(0.5)
        title_paragraph.font.bold = True
    
    # 添加图片
    if i > 0:  # 第一页不添加图片
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(15), height=Inches(7))

# 保存PPT
ppt_path = '/Users/tianyong/skills/2026一季度公司整体运营总览.pptx'
prs.save(ppt_path)
print(f'PPT生成完成: {ppt_path}')